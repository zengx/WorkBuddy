#!/usr/bin/env python3
"""Extract text (Markdown) + images from .docx / .pdf / .pptx for sheet-generation reference.

Pipeline:
1. Extract text as Markdown (mammoth for docx, pymupdf for pdf, python-pptx
   for pptx — slide text frames, tables, charts, speaker notes).
2. Extract original images to images/ (high-res, retained for traceability).
3. Generate thumbnails/ (small PNG previews; sheet-agent can Read these as
   visual fallback when content.md text alone is insufficient — e.g. scanned
   PDFs, table screenshots, or slide layouts dominated by images).

Default output is a subfolder under the project's resources/ tree:
    <cwd>/resources/extracted/<stem>/{content.md, images/, thumbnails/}

The excel-generation skill explicitly overrides the output via -o to place
artifacts under the ref directory:
    <ref_dir>/reference/<stem>/{content.md, images/, thumbnails/}

Multiple reference docs can coexist without filename collisions.

Thumbnails are uniformly PNG to preserve text legibility (no JPEG artifacts
on text edges). Token cost in vision models depends on resolution, not file
size, so PNG is preferred for agent viewing.

Safety: this script never performs wildcard deletion. It writes into a new
output directory and overwrites only explicit files inside it.

Dependencies:
    - mammoth           (docx → HTML with image extraction)
    - pymupdf           (PDF text + embedded image extraction)
    - python-pptx       (PPTX text frames / tables / images / charts / notes)
    - Pillow            (thumbnail generation)
    - markdownify       (HTML → Markdown conversion)

Usage:
    python3 scripts/extract.py <input_file> [-o output_dir] [--thumb-width 600]

Examples:
    # Default: writes to ./resources/extracted/report/
    python3 scripts/extract.py ~/Desktop/report.docx

    # Explicit output directory (excel-generation skill uses this form)
    python3 scripts/extract.py paper.pdf -o /path/to/ref/reference/paper

    # Slide deck as content reference
    python3 scripts/extract.py q1_review.pptx -o /path/to/ref/reference/q1_review

    # Higher-res thumbnails for dense screenshots
    python3 scripts/extract.py document.docx --thumb-width 800
"""

from __future__ import annotations

import argparse
import subprocess
import sys
from pathlib import Path


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _die(message: str, code: int = 1) -> None:
    print(f"[extract] ERROR: {message}", file=sys.stderr)
    raise SystemExit(code)


def _auto_install(import_name: str, pip_name: str | None = None) -> None:
    """Try importing a package; auto-install via pip if missing."""
    pip_name = pip_name or import_name
    try:
        __import__(import_name)
    except ImportError:
        print(f"[extract] {pip_name} not found, installing ...", file=sys.stderr)
        try:
            subprocess.check_call(
                [sys.executable, "-m", "pip", "install", "--quiet", pip_name]
            )
        except subprocess.CalledProcessError as e:
            _die(f"failed to install {pip_name}: {e}", code=4)
        # Verify after install
        try:
            __import__(import_name)
        except ImportError as e:
            _die(f"{pip_name} still unavailable after install: {e}", code=4)


def _check_deps(
    need_pymupdf: bool = False,
    need_mammoth: bool = False,
    need_pptx: bool = False,
) -> None:
    _auto_install("PIL", "Pillow")

    if need_mammoth:
        _auto_install("mammoth")
        _auto_install("markdownify")

    if need_pymupdf:
        _auto_install("fitz", "pymupdf")

    if need_pptx:
        _auto_install("pptx", "python-pptx")


# ---------------------------------------------------------------------------
# Thumbnail generation
# ---------------------------------------------------------------------------


def _make_thumbnail(
    src: Path,
    dst_dir: Path,
    max_width: int = 600,
) -> Path:
    """Create a PNG thumbnail (uniform format for legible text)."""
    from PIL import Image

    img = Image.open(src)
    # Preserve transparency where present; convert palette/CMYK to RGB
    if img.mode in ("P", "CMYK"):
        img = img.convert("RGB")

    # Resize proportionally
    if img.width > max_width:
        ratio = max_width / img.width
        new_size = (max_width, int(img.height * ratio))
        img = img.resize(new_size, Image.LANCZOS)

    out_path = dst_dir / (src.stem + ".png")
    img.save(out_path, "PNG", optimize=True)
    return out_path


def _generate_thumbnails(
    images_dir: Path,
    thumbs_dir: Path,
    max_width: int,
) -> list[Path]:
    """Generate PNG thumbnails for all images in a directory."""
    thumbs_dir.mkdir(parents=True, exist_ok=True)
    results = []

    image_exts = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp"}
    images = sorted(
        f for f in images_dir.iterdir() if f.suffix.lower() in image_exts
    )
    if not images:
        return results

    for img_path in images:
        try:
            thumb = _make_thumbnail(img_path, thumbs_dir, max_width)
            results.append(thumb)
        except Exception as e:
            print(f"  [extract] WARNING: failed to thumbnail {img_path.name}: {e}")

    return results


# ---------------------------------------------------------------------------
# DOCX extraction
# ---------------------------------------------------------------------------


def _extract_docx(docx_path: Path, out_dir: Path) -> tuple[Path, int]:
    """Extract text via mammoth (file-path image refs) + markdownify."""
    import mammoth
    from markdownify import markdownify

    images_dir = out_dir / "images"
    images_dir.mkdir(parents=True, exist_ok=True)

    img_counter = 0

    @mammoth.images.img_element
    def save_image(image):
        nonlocal img_counter
        img_counter += 1
        ext_map = {
            "image/png": ".png",
            "image/jpeg": ".jpg",
            "image/gif": ".gif",
            "image/bmp": ".bmp",
            "image/tiff": ".tiff",
            "image/webp": ".webp",
        }
        ext = ext_map.get(image.content_type, ".png")
        fname = f"image{img_counter}{ext}"
        with image.open() as img_bytes:
            (images_dir / fname).write_bytes(img_bytes.read())
        return {"src": f"images/{fname}"}

    with open(docx_path, "rb") as f:
        result = mammoth.convert_to_html(f, convert_image=save_image)

    html_content = result.value

    # Convert HTML → Markdown
    md_content = markdownify(html_content, heading_style="ATX", strip=["script", "style"])
    md_content = md_content.strip()

    md_file = out_dir / "content.md"
    md_file.write_text(md_content, encoding="utf-8")

    return md_file, img_counter


# ---------------------------------------------------------------------------
# PDF extraction
# ---------------------------------------------------------------------------


def _extract_pdf(pdf_path: Path, out_dir: Path) -> tuple[Path, int]:
    """Extract text + embedded images from PDF via pymupdf."""
    import fitz

    images_dir = out_dir / "images"
    images_dir.mkdir(parents=True, exist_ok=True)

    md_parts = []
    img_count = 0

    with fitz.open(pdf_path) as doc:
        for page_idx in range(len(doc)):
            page = doc[page_idx]
            page_num = page_idx + 1

            md_parts.append(f"## Page {page_num}\n")

            text = page.get_text("text").strip()
            if text:
                md_parts.append(text)
                md_parts.append("")

            image_list = page.get_images(full=True)
            for img_idx, img_info in enumerate(image_list, 1):
                xref = img_info[0]
                try:
                    pix = fitz.Pixmap(doc, xref)
                    # Only convert CMYK (n - alpha > 3) to RGB. RGB / RGBA /
                    # grayscale are saved as-is — PNG supports all of them
                    # and we want to preserve transparency.
                    if pix.n - pix.alpha > 3:
                        pix = fitz.Pixmap(fitz.csRGB, pix)

                    img_name = f"page{page_num:02d}_img{img_idx:02d}.png"
                    img_path = images_dir / img_name
                    pix.save(str(img_path))
                    img_count += 1

                    md_parts.append(f"![](images/{img_name})")
                    md_parts.append("")
                except Exception as e:
                    print(f"  [extract] WARNING: page {page_num} image {img_idx}: {e}")

    md_content = "\n".join(md_parts)
    md_file = out_dir / "content.md"
    md_file.write_text(md_content, encoding="utf-8")

    return md_file, img_count


# ---------------------------------------------------------------------------
# PPTX extraction
# ---------------------------------------------------------------------------


def _pptx_table_to_md(table) -> str:
    """Render a python-pptx table as a markdown table.

    First row is treated as the header. Cells with line breaks are flattened
    with `<br>` so the markdown renders as a single row each. Pipes inside
    cell text are escaped to avoid breaking the table.
    """

    def cell_text(cell) -> str:
        # Cells contain a text_frame; .text already concatenates paragraphs
        # with newlines. Flatten + escape pipes for markdown safety.
        return cell.text.replace("|", "\\|").replace("\n", "<br>").strip()

    rows = list(table.rows)
    if not rows:
        return ""

    header = [cell_text(c) for c in rows[0].cells]
    body = [[cell_text(c) for c in r.cells] for r in rows[1:]]

    # Pad/truncate body rows to header width for safety
    width = len(header)
    norm_body = [(row + [""] * width)[:width] for row in body]

    lines = ["| " + " | ".join(header) + " |"]
    lines.append("|" + "|".join(["---"] * width) + "|")
    for row in norm_body:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def _pptx_chart_to_md(chart) -> str:
    """Best-effort markdown summary of a chart: type + categories + series.

    python-pptx exposes plain category labels and numeric series values for
    common chart types (bar/line/pie/area/etc.). For uncommon types it may
    raise; in that case we fall back to a one-line note.
    """
    try:
        chart_type = str(chart.chart_type).split(".")[-1].rstrip(">")
        plot = chart.plots[0]
        categories = list(plot.categories)
        series_blocks = []
        for s in plot.series:
            try:
                values = list(s.values)
            except Exception:
                values = []
            series_blocks.append((s.name or "<unnamed>", values))

        lines = [f"**Chart** ({chart_type})", ""]
        if categories or series_blocks:
            header = ["category"] + [name for name, _ in series_blocks]
            lines.append("| " + " | ".join(str(h) for h in header) + " |")
            lines.append("|" + "|".join(["---"] * len(header)) + "|")
            for i, cat in enumerate(categories):
                row = [str(cat)]
                for _, values in series_blocks:
                    row.append(str(values[i]) if i < len(values) else "")
                lines.append("| " + " | ".join(row) + " |")
        return "\n".join(lines)
    except Exception as e:
        return f"**Chart** (extraction failed: {e})"


def _extract_pptx(pptx_path: Path, out_dir: Path) -> tuple[Path, int]:
    """Extract text frames, tables, charts, images, and notes from a .pptx.

    Layout in content.md:
        ## Slide N
        <text frames>
        ### Tables
        <markdown tables>
        ### Charts
        <markdown chart summaries>
        ![](images/slideN_imgM.ext)
        > **Speaker Notes:** ...
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    images_dir = out_dir / "images"
    images_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(pptx_path)
    md_parts: list[str] = []
    img_count = 0

    for slide_idx, slide in enumerate(prs.slides, 1):
        md_parts.append(f"## Slide {slide_idx}\n")

        text_blocks: list[str] = []
        table_blocks: list[str] = []
        chart_blocks: list[str] = []
        image_refs: list[str] = []

        # Walk top-level shapes; nested shapes inside groups are recursed via
        # an explicit stack so we don't miss text inside grouped layouts.
        stack = list(slide.shapes)
        while stack:
            shape = stack.pop(0)

            # Recurse into groups
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                stack[0:0] = list(shape.shapes)
                continue

            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt:
                    text_blocks.append(txt)

            if shape.has_table:
                md_table = _pptx_table_to_md(shape.table)
                if md_table:
                    table_blocks.append(md_table)

            if shape.has_chart:
                chart_blocks.append(_pptx_chart_to_md(shape.chart))

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = shape.image
                    ext = (img.ext or "png").lstrip(".").lower()
                    img_count += 1
                    fname = f"slide{slide_idx:02d}_img{img_count:02d}.{ext}"
                    (images_dir / fname).write_bytes(img.blob)
                    image_refs.append(f"![](images/{fname})")
                except Exception as e:
                    print(
                        f"  [extract] WARNING: slide {slide_idx} image: {e}",
                        file=sys.stderr,
                    )

        if text_blocks:
            md_parts.append("\n\n".join(text_blocks))
            md_parts.append("")

        if table_blocks:
            md_parts.append("### Tables\n")
            md_parts.append("\n\n".join(table_blocks))
            md_parts.append("")

        if chart_blocks:
            md_parts.append("### Charts\n")
            md_parts.append("\n\n".join(chart_blocks))
            md_parts.append("")

        if image_refs:
            md_parts.append("\n".join(image_refs))
            md_parts.append("")

        # Speaker notes (skip the empty placeholder pptx attaches by default)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                quoted = "\n".join(f"> {line}" for line in notes.splitlines())
                md_parts.append(f"> **Speaker Notes:**\n{quoted}")
                md_parts.append("")

    md_content = "\n".join(md_parts).rstrip() + "\n"
    md_file = out_dir / "content.md"
    md_file.write_text(md_content, encoding="utf-8")

    return md_file, img_count


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Extract text + images from .docx/.pdf/.pptx for sheet generation reference"
    )
    parser.add_argument("input", help="Input file (.docx, .pdf, or .pptx)")
    parser.add_argument(
        "-o", "--out", default=None,
        help="Output directory (default: <cwd>/resources/extracted/<stem>/)",
    )
    parser.add_argument(
        "--thumb-width", type=int, default=600,
        help="Max thumbnail width in px (default: 600, balances legibility and tokens)",
    )
    args = parser.parse_args()

    if args.thumb_width < 50:
        _die(f"--thumb-width must be at least 50, got {args.thumb_width}")

    input_path = Path(args.input).expanduser().resolve()
    if not input_path.exists():
        _die(f"file not found: {input_path}")

    suffix = input_path.suffix.lower()
    if suffix not in (".docx", ".pdf", ".pptx"):
        _die(f"unsupported format '{suffix}'. Supported: .docx, .pdf, .pptx")

    _check_deps(
        need_pymupdf=(suffix == ".pdf"),
        need_mammoth=(suffix == ".docx"),
        need_pptx=(suffix == ".pptx"),
    )

    out_dir = (
        Path(args.out).expanduser().resolve()
        if args.out
        else Path.cwd() / "resources" / "extracted" / input_path.stem
    )
    out_dir.mkdir(parents=True, exist_ok=True)

    print(f"[extract] Input:  {input_path}")
    print(f"[extract] Output: {out_dir}")
    print()

    # Extract
    if suffix == ".docx":
        print("[extract] [1/2] Extracting from DOCX (mammoth)...")
        md_file, img_count = _extract_docx(input_path, out_dir)
    elif suffix == ".pdf":
        print("[extract] [1/2] Extracting from PDF (pymupdf)...")
        md_file, img_count = _extract_pdf(input_path, out_dir)
    else:
        print("[extract] [1/2] Extracting from PPTX (python-pptx)...")
        md_file, img_count = _extract_pptx(input_path, out_dir)

    print(f"[extract]   Text → {md_file.name}")
    print(f"[extract]   Images: {img_count} found")

    # Thumbnails
    images_dir = out_dir / "images"
    thumbs_dir = out_dir / "thumbnails"

    if img_count > 0 and images_dir.exists():
        print(f"\n[extract] [2/2] Generating PNG thumbnails (max {args.thumb_width}px)...")
        thumbs = _generate_thumbnails(images_dir, thumbs_dir, args.thumb_width)
        print(f"[extract]   Thumbnails: {len(thumbs)} PNG files")
    else:
        print("\n[extract] [2/2] No images to thumbnail, skipping.")

    # Summary
    print()
    print("=" * 60)
    print("[extract] Done! Output structure:")
    print(f"  {out_dir}/")
    print(f"    content.md        ← text + image refs (agent reads this)")
    print(f"    images/           ← originals ({img_count} files, for slides)")
    if thumbs_dir.exists():
        print(f"    thumbnails/       ← previews (agent views these)")
    print("=" * 60)


if __name__ == "__main__":
    main()